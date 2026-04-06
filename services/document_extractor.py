"""
Document Extractor Service
==========================
Centralised text extraction for multiple file formats.

Supported formats:
  - PDF  (.pdf)   — via PyPDF2
  - DOCX (.docx)  — via python-docx
  - PPTX (.pptx)  — via python-pptx
  - TXT  (.txt)   — plain read
"""

import os
import csv

import PyPDF2
from werkzeug.utils import secure_filename


ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16 MB


# ── Public API ────────────────────────────────────────────────────────────────

class DocumentExtractor:
    """Service for extracting text from uploaded documents."""

    @staticmethod
    def allowed_file(filename: str) -> bool:
        """Check whether a filename has an allowed extension."""
        return '.' in filename and \
               filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

    @staticmethod
    def extract_text(file_path: str) -> str:
        """
        Dispatch to the correct extraction method based on extension.

        Returns:
            Extracted text as a single string.

        Raises:
            Exception: On unsupported format or extraction failure.
        """
        ext = os.path.splitext(file_path)[1].lower()
        extractors = {
            '.pdf':  DocumentExtractor._extract_pdf,
            '.docx': DocumentExtractor._extract_docx,
            '.pptx': DocumentExtractor._extract_pptx,
            '.txt':  DocumentExtractor._extract_txt,
        }

        extractor = extractors.get(ext)
        if extractor is None:
            raise Exception(
                f"Unsupported file type: {ext}. "
                f"Allowed types: {', '.join(ALLOWED_EXTENSIONS)}"
            )

        text = extractor(file_path)
        if not text or not text.strip():
            raise Exception("No text content could be extracted from the file")
        return text

    @staticmethod
    def save_uploaded_file(file, upload_folder: str = 'uploads') -> str:
        """Save an uploaded FileStorage object to disk and return the path."""
        os.makedirs(upload_folder, exist_ok=True)

        filename = secure_filename(file.filename)
        file_path = os.path.join(upload_folder, filename)

        # Deduplicate filename
        base_name, extension = os.path.splitext(filename)
        counter = 1
        while os.path.exists(file_path):
            filename = f"{base_name}_{counter}{extension}"
            file_path = os.path.join(upload_folder, filename)
            counter += 1

        file.save(file_path)
        return file_path

    @staticmethod
    def process_upload(file):
        """
        Complete workflow: validate → save → extract text.

        Args:
            file: Flask FileStorage object.

        Returns:
            Tuple of (file_path, extracted_text, source_type).
        """
        if not file or file.filename == '':
            raise Exception("No file provided")

        if not DocumentExtractor.allowed_file(file.filename):
            allowed = ', '.join(f'.{e}' for e in sorted(ALLOWED_EXTENSIONS))
            raise Exception(
                f"Unsupported file type. Allowed formats: {allowed}"
            )

        file_path = DocumentExtractor.save_uploaded_file(file)

        try:
            text = DocumentExtractor.extract_text(file_path)
            source_type = os.path.splitext(file.filename)[1].lstrip('.').lower()
            return file_path, text, source_type
        except Exception:
            # Clean up on failure
            if os.path.exists(file_path):
                os.remove(file_path)
            raise


# ── Private Extractors ────────────────────────────────────────────────────────

    @staticmethod
    def _extract_pdf(path: str) -> str:
        """Extract text from a PDF using PyPDF2."""
        try:
            pages = []
            with open(path, 'rb') as fh:
                reader = PyPDF2.PdfReader(fh)
                if reader.is_encrypted:
                    raise Exception("PDF is encrypted and cannot be read")
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        pages.append(text)
            full = '\n\n'.join(pages)
            return ' '.join(full.split())  # collapse whitespace
        except Exception as e:
            raise Exception(f"Error extracting text from PDF: {e}")

    @staticmethod
    def _extract_docx(path: str) -> str:
        """Extract text from a Word document (.docx)."""
        try:
            from docx import Document
            doc = Document(path)

            parts = []
            # Paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)

            # Tables
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(' | '.join(cells))

            return '\n'.join(parts)
        except Exception as e:
            raise Exception(f"Error extracting text from DOCX: {e}")

    @staticmethod
    def _extract_pptx(path: str) -> str:
        """Extract text from a PowerPoint presentation (.pptx)."""
        try:
            from pptx import Presentation
            prs = Presentation(path)

            parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_texts.append(text)

                if slide_texts:
                    parts.append(f"--- Slide {slide_num} ---")
                    parts.extend(slide_texts)

            return '\n'.join(parts)
        except Exception as e:
            raise Exception(f"Error extracting text from PPTX: {e}")

    @staticmethod
    def _extract_txt(path: str) -> str:
        """Read plain text with encoding fallback."""
        for encoding in ('utf-8', 'latin-1', 'cp1252'):
            try:
                with open(path, 'r', encoding=encoding) as fh:
                    return fh.read()
            except UnicodeDecodeError:
                continue
        raise Exception("Could not decode text file with any supported encoding")
